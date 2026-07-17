"""Pack 14: PPTX export builder — create, add slides, finalize to tenant exports. Returns file_id for download."""

from __future__ import annotations

import uuid
from pathlib import Path
from typing import Any, Dict, List

from hg_core.docs.paths import get_exports_root

_pptx_buffers: Dict[str, Any] = {}


def pptx_create(title: str, tenant_id: str) -> str:
    """Create a new PPTX buffer. Returns doc_id for add_slide / finalize."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx required; pip install python-pptx")
    doc_id = str(uuid.uuid4())
    prs = Presentation()
    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    _pptx_buffers[doc_id] = {"prs": prs, "tenant_id": tenant_id}
    return doc_id


def pptx_add_slide(doc_id: str, content: str) -> None:
    """Add a slide with the given text content."""
    buf = _pptx_buffers.get(doc_id)
    if not buf:
        raise ValueError("doc_id not found")
    try:
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx required")
    prs = buf["prs"]
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
    p = tx.text_frame.paragraphs[0]
    p.text = content or ""
    p.font.size = Pt(18)
    tx.text_frame.word_wrap = True


def pptx_finalize(doc_id: str, filename: str) -> str:
    """Write PPTX to tenant exports dir. Returns file_id for GET /v1/files/{file_id}/download."""
    buf = _pptx_buffers.pop(doc_id, None)
    if not buf:
        raise ValueError("doc_id not found")
    tenant_id = buf["tenant_id"]
    export_root = get_exports_root(tenant_id)
    export_root.mkdir(parents=True, exist_ok=True)
    file_id = str(uuid.uuid4())
    safe_name = (filename or "export").strip().replace("..", "") or "export"
    if not safe_name.lower().endswith(".pptx"):
        safe_name += ".pptx"
    path = export_root / f"{file_id}.pptx"
    buf["prs"].save(str(path))
    return file_id
