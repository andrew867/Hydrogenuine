"""Office Agent tool pack: docx, xlsx, pptx, pdf. Phase 11. Idempotency per call (path + op)."""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any, Dict, List, Optional

from .tool_router import ToolCall


def _get_workspace_root() -> Optional[Path]:
    try:
        from hg_lib.config import get_workspace_root
        return get_workspace_root()
    except Exception:
        return None


def _get(call: ToolCall, key: str, default: Any = None) -> Any:
    return call.args.get(key, default)


def _safe_path(path_arg: str, workspace: Optional[Path]) -> Optional[Path]:
    if not path_arg or not isinstance(path_arg, str):
        return None
    if workspace is None:
        return None
    p = Path(path_arg)
    if p.is_absolute():
        resolved = p
    else:
        resolved = (workspace / path_arg).resolve()
    try:
        resolved.relative_to(workspace.resolve())
    except ValueError:
        return None
    return resolved


def _path_and_workspace(call: ToolCall) -> tuple[Optional[Path], Optional[Path]]:
    path_arg = _get(call, "path")
    if not path_arg:
        return None, None
    ws = _get(call, "workspace")
    if ws is not None:
        ws = Path(ws) if isinstance(ws, str) else None
    else:
        ws = _get_workspace_root()
    path = _safe_path(str(path_arg), ws)
    return path, ws


# --- docx ---

def handler_office_docx_read(call: ToolCall) -> Dict[str, Any]:
    """Read .docx file; return paragraphs and tables as structured data."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.docx.read"}
    if not path.exists() or not path.is_file():
        return {"ok": False, "error": f"file not found: {path}", "action": "office.docx.read"}
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs]
        tables = []
        for t in doc.tables:
            rows = [[c.text for c in row.cells] for row in t.rows]
            tables.append(rows)
        return {
            "ok": True,
            "data": {"path": str(path), "paragraphs": paragraphs, "tables": tables},
            "action": "office.docx.read",
        }
    except ImportError:
        return {"ok": False, "error": "python-docx not installed; pip install python-docx", "action": "office.docx.read"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.docx.read"}


def handler_office_docx_write(call: ToolCall) -> Dict[str, Any]:
    """Write content to .docx file. args: path, content (str or list of paragraphs)."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.docx.write"}
    content = _get(call, "content")
    if content is None:
        return {"ok": False, "error": "content required", "action": "office.docx.write"}
    try:
        from docx import Document
        doc = Document()
        if isinstance(content, list):
            for p in content:
                doc.add_paragraph(p if isinstance(p, str) else str(p))
        else:
            doc.add_paragraph(str(content))
        path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(path))
        return {"ok": True, "data": {"path": str(path)}, "action": "office.docx.write"}
    except ImportError:
        return {"ok": False, "error": "python-docx not installed; pip install python-docx", "action": "office.docx.write"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.docx.write"}


# --- xlsx ---

def handler_office_xlsx_read(call: ToolCall) -> Dict[str, Any]:
    """Read .xlsx file; return sheet names and sheet data (list of rows)."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.xlsx.read"}
    if not path.exists() or not path.is_file():
        return {"ok": False, "error": f"file not found: {path}", "action": "office.xlsx.read"}
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        sheet_names = list(wb.sheetnames)
        sheets = {}
        for name in sheet_names:
            sh = wb[name]
            rows = [[str(c.value) if c.value is not None else "" for c in row] for row in sh.iter_rows()]
            sheets[name] = rows
        wb.close()
        return {
            "ok": True,
            "data": {"path": str(path), "sheets": sheets, "sheet_names": sheet_names},
            "action": "office.xlsx.read",
        }
    except ImportError:
        return {"ok": False, "error": "openpyxl not installed; pip install openpyxl", "action": "office.xlsx.read"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.xlsx.read"}


def handler_office_xlsx_write(call: ToolCall) -> Dict[str, Any]:
    """Write data to .xlsx file. args: path, data (list of rows or dict sheet_name -> rows), sheet_name (optional)."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.xlsx.write"}
    data = _get(call, "data")
    if data is None:
        return {"ok": False, "error": "data required", "action": "office.xlsx.write"}
    try:
        import openpyxl
        wb = openpyxl.Workbook()
        if isinstance(data, dict):
            for sheet_name, rows in data.items():
                ws = wb.create_sheet(title=str(sheet_name)[:31])
                for r in rows:
                    ws.append(r if isinstance(r, (list, tuple)) else [r])
        else:
            rows = data if isinstance(data, list) else [[data]]
            for r in rows:
                wb.active.append(r if isinstance(r, (list, tuple)) else [r])
        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(path))
        return {"ok": True, "data": {"path": str(path)}, "action": "office.xlsx.write"}
    except ImportError:
        return {"ok": False, "error": "openpyxl not installed; pip install openpyxl", "action": "office.xlsx.write"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.xlsx.write"}


# --- pptx ---

def handler_office_pptx_read(call: ToolCall) -> Dict[str, Any]:
    """Read .pptx file; return slide text and shape text."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.pptx.read"}
    if not path.exists() or not path.is_file():
        return {"ok": False, "error": f"file not found: {path}", "action": "office.pptx.read"}
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            slides.append({"slide_index": i, "texts": texts})
        return {
            "ok": True,
            "data": {"path": str(path), "slides": slides},
            "action": "office.pptx.read",
        }
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed; pip install python-pptx", "action": "office.pptx.read"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.pptx.read"}


def handler_office_pptx_write(call: ToolCall) -> Dict[str, Any]:
    """Write content to .pptx file. args: path, content (str or list of strings per slide)."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.pptx.write"}
    content = _get(call, "content")
    if content is None:
        return {"ok": False, "error": "content required", "action": "office.pptx.write"}
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        items = content if isinstance(content, list) else [content]
        for i, item in enumerate(items):
            text = item if isinstance(item, str) else str(item)
            slide_layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(slide_layout)
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            p = tx.text_frame.paragraphs[0]
            p.text = text
            p.font.size = Pt(18)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return {"ok": True, "data": {"path": str(path)}, "action": "office.pptx.write"}
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed; pip install python-pptx", "action": "office.pptx.write"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.pptx.write"}


# --- pdf ---

def handler_office_pdf_read(call: ToolCall) -> Dict[str, Any]:
    """Read .pdf file; extract text per page."""
    path, _ = _path_and_workspace(call)
    if path is None:
        return {"ok": False, "error": "path required or invalid", "action": "office.pdf.read"}
    if not path.exists() or not path.is_file():
        return {"ok": False, "error": f"file not found: {path}", "action": "office.pdf.read"}
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({"page": i + 1, "text": text})
        return {
            "ok": True,
            "data": {"path": str(path), "pages": pages, "num_pages": len(reader.pages)},
            "action": "office.pdf.read",
        }
    except ImportError:
        return {"ok": False, "error": "pypdf not installed; pip install pypdf", "action": "office.pdf.read"}
    except Exception as e:
        return {"ok": False, "error": str(e), "action": "office.pdf.read"}


def idempotency_key_office(path: str, op: str) -> str:
    """Stable idempotency key for office tools (min 8 chars)."""
    h = hashlib.sha256(f"office:{op}:{path}".encode("utf-8", errors="replace")).hexdigest()
    return f"office-{op}-{h[:20]}"
